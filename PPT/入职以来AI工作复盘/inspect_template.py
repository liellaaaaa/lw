#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Dump a PPTX template's per-slide structure: layout, shapes, text, fonts, colors.
Used to reverse-engineer the visual design so we can replicate it in DESIGN.md.
"""
from pptx import Presentation
from pptx.util import Emu
import sys, os

SRC = r"C:/Users/windows/Desktop/月度经营分析会汇报模版：PBC执行进度.pptx"

def emu_cm(v):
    return round(v / 914400 * 2.54, 2) if v is not None else None

def emu_px(v):
    return round(v / 9525) if v is not None else None

def color_of(el):
    # el is a ColorFormat (font.color / fill.fore_color)
    try:
        if el is None:
            return None
        if el.type is None:
            return None
        if el.type == 1:  # RGB
            return str(el.rgb)
        else:
            return f"theme:{el.theme_color}"
    except Exception as e:
        return f"err:{e}"

def walk_shape(sh, depth=1, out=None):
    if out is None:
        out = []
    ind = "  " * depth
    st = sh.shape_type
    left = emu_cm(sh.left) if sh.left is not None else None
    top = emu_cm(sh.top) if sh.top is not None else None
    w = emu_cm(sh.width) if sh.width is not None else None
    h = emu_cm(sh.height) if sh.height is not None else None
    line = f"{ind}{st} name={sh.name!r} pos=({left},{top}) size=({w}x{h})cm"
    # fill
    try:
        if sh.fill.type is not None:
            fc = color_of(sh.fill.fore_color) if sh.fill.type == 1 else None
            line += f" fill={sh.fill.type}:{fc}"
    except Exception:
        pass
    # line color
    try:
        if sh.line.color and sh.line.color.type is not None:
            line += f" line={color_of(sh.line.color)}"
    except Exception:
        pass
    out.append(line)
    # text
    if sh.has_text_frame:
        for pi, p in enumerate(sh.text_frame.paragraphs):
            txt = "".join(r.text for r in p.runs)
            if not txt.strip():
                continue
            runs_info = []
            for r in p.runs:
                f = r.font
                runs_info.append(f"{r.text!r}[sz={f.size.pt if f.size else None},b={f.bold},i={f.italic},name={f.name},rgb={color_of(f.color)}]")
            out.append(f"{ind}  P{pi}: {txt[:80]!r}")
            for ri in runs_info[:6]:
                out.append(f"{ind}    run {ri}")
    # table
    if sh.has_table:
        tbl = sh.table
        out.append(f"{ind}  TABLE rows={len(tbl.rows)} cols={len(tbl.columns)}")
        for ri, row in enumerate(tbl.rows):
            cells = [c.text[:18] for c in row.cells]
            out.append(f"{ind}    R{ri}: {cells}")
            # cell fill
            for ci, c in enumerate(row.cells):
                try:
                    if c.fill.type == 1:
                        out.append(f"{ind}      cell[{ri},{ci}] fill={color_of(c.fill.fore_color)}")
                except Exception:
                    pass
    # picture
    if st == 13:  # PICTURE
        out.append(f"{ind}  [PICTURE]")
    # group / children
    if sh.shape_type == 6:  # GROUP
        for child in sh.shapes:
            walk_shape(child, depth + 1, out)
    return out

def main():
    prs = Presentation(SRC)
    W = emu_cm(prs.slide_width); H = emu_cm(prs.slide_height)
    print(f"SLIDE SIZE: {W} x {H} cm  ({emu_px(prs.slide_width)}x{emu_px(prs.slide_height)} px)")
    print(f"SLIDES: {len(prs.slides)}")
    # theme colors
    try:
        from pptx.oxml.ns import qn
        theme = prs.slide_masters[0].element
        for tag in ['dk1','lt1','dk2','lt2','accent1','accent2','accent3','accent4','accent5','accent6']:
            els = theme.findall(f'.//{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}')
            for e in els:
                srgb = e.find(qn('a:srgbClr'))
                if srgb is not None:
                    print(f"THEME {tag} = #{srgb.get('val')}")
    except Exception as e:
        print(f"THEME err {e}")
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n===== SLIDE {i} =====")
        try:
            bg = slide.background
            if bg.fill.type is not None:
                print(f"  BG fill={bg.fill.type} {color_of(bg.fill.fore_color)}")
        except Exception as e:
            print(f"  BG err {e}")
        # layout name
        try:
            print(f"  LAYOUT: {slide.slide_layout.name}")
        except Exception:
            pass
        for sh in slide.shapes:
            for ln in walk_shape(sh, 1):
                print("  " + ln)

if __name__ == "__main__":
    main()
